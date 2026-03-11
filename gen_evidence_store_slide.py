"""Generate a single PPT slide: Why Train-Only Evidence Store Is Sufficient."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Background ───────────────────────────────────────────────────
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# ── Helper ───────────────────────────────────────────────────────
def add_text_box(left, top, width, height, text, font_size=14,
                 bold=False, color=RGBColor(0x33, 0x33, 0x33),
                 align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_para(tf, text, font_size=14, bold=False,
             color=RGBColor(0x33, 0x33, 0x33), bullet=False,
             space_before=Pt(4), indent=0, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.space_before = space_before
    if indent:
        p.level = indent
    return p

def add_rounded_rect(left, top, width, height, fill_rgb):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

# ── Title bar ────────────────────────────────────────────────────
title_bar = add_rounded_rect(0.3, 0.2, 12.7, 0.7, RGBColor(0x2C, 0x3E, 0x50))
add_text_box(0.5, 0.25, 12, 0.6,
             "Why Train-Only Evidence Store Is Sufficient for RAG Retrieval",
             font_size=24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.LEFT)

# ── Subtitle ─────────────────────────────────────────────────────
add_text_box(0.5, 1.0, 12, 0.4,
             "Advisor Question: \"Why not use the full dataset for the evidence store (as in typical RAG)?\"",
             font_size=15, bold=False, color=RGBColor(0x66, 0x66, 0x66))

# ── LEFT COLUMN: Design Rationale ────────────────────────────────
left_header_bg = add_rounded_rect(0.4, 1.55, 6.1, 0.45, RGBColor(0x34, 0x95, 0xDB))
add_text_box(0.55, 1.58, 5.8, 0.4,
             "Design Rationale", font_size=16, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))

tf = add_text_box(0.55, 2.1, 5.9, 3.5,
                  "1. Prevents data leakage  --  Using only train-set logs "
                  "ensures the retriever never sees test-set anomalies during "
                  "evidence construction.",
                  font_size=12.5, color=RGBColor(0x33, 0x33, 0x33))
add_para(tf, "", font_size=6)
add_para(tf, "2. Simulates real deployment  --  In production, only historical "
         "(past) logs are available when a new anomaly arrives. Train-only "
         "mirrors this temporal split.", font_size=12.5)
add_para(tf, "", font_size=6)
add_para(tf, "3. Controls for confounding  --  If test anomalies appear in the "
         "evidence store, the LLM could retrieve near-identical sessions, "
         "inflating retrieval metrics artificially.", font_size=12.5)
add_para(tf, "", font_size=6)
add_para(tf, "4. Standard ML practice  --  Train/test separation is fundamental. "
         "RAG's retrieval corpus is part of the model's knowledge base and "
         "must respect the same split.", font_size=12.5)

# ── RIGHT COLUMN: Empirical Evidence ─────────────────────────────
right_header_bg = add_rounded_rect(6.8, 1.55, 6.1, 0.45, RGBColor(0x27, 0xAE, 0x60))
add_text_box(6.95, 1.58, 5.8, 0.4,
             "Empirical Evidence (RQ2 Phase 1)", font_size=16, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))

# ── Overlap table ────────────────────────────────────────────────
from pptx.util import Cm

cols = 3
rows = 8
tbl_left = Inches(7.0)
tbl_top = Inches(2.15)
tbl_w = Inches(5.7)
tbl_h = Inches(2.7)

table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h)
table = table_shape.table

# column widths
table.columns[0].width = Inches(2.8)
table.columns[1].width = Inches(1.45)
table.columns[2].width = Inches(1.45)

data = [
    ["Metric", "BGL", "HDFS"],
    ["Train-only corpus size", "332,356", "402,542"],
    ["Whole-dataset corpus size", "474,796", "575,061"],
    ["Anomaly overlap rate (E1-E4)", "88.25%", "86.75%"],
    ["Top-1 anomaly match rate", "85.0%", "86.0%"],
    ["Perfect overlap (4/4 match)", "70.0%", "70.0%"],
    ["New docs in anomaly top-4", "0.47 / 4", "0.31 / 4"],
    ["Sessions with 0 new docs", "70%", "76%"],
]

header_fill = RGBColor(0x2C, 0x3E, 0x50)
alt_fill = RGBColor(0xF0, 0xF4, 0xF7)

for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = data[r][c]
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.name = "Calibri"
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        else:
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            if c == 0:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = True

# ── Bottom section: Key Findings ──────────────────────────────────
findings_bg = add_rounded_rect(0.4, 5.1, 12.5, 2.0, RGBColor(0xF9, 0xF9, 0xF9))
findings_border = add_rounded_rect(0.4, 5.1, 12.5, 2.0, RGBColor(0xDD, 0xDD, 0xDD))
findings_border.fill.background()

findings_header = add_rounded_rect(0.4, 5.1, 12.5, 0.4, RGBColor(0xE8, 0x4C, 0x3D))
add_text_box(0.55, 5.12, 5, 0.35,
             "Key Findings", font_size=15, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))

tf2 = add_text_box(0.6, 5.6, 12.1, 1.4,
    "1.  88% / 87% anomaly overlap between train-only and whole-dataset retrieval "
    "(BGL / HDFS)  --  adding test data changes < 0.5 docs per query on average.",
    font_size=12.5, color=RGBColor(0x33, 0x33, 0x33))
add_para(tf2, "", font_size=4)
add_para(tf2,
    "2.  70% of test sessions get identical anomaly evidence (4/4 perfect match) "
    "regardless of corpus.  The marginal gain from the extra 40% data is negligible.",
    font_size=12.5)
add_para(tf2, "", font_size=4)
add_para(tf2,
    "3.  Context Precision@4 > 0.95 and Grounding Breadth ~66% already achieved "
    "with train-only corpus  --  retrieval quality is saturated before adding test data.",
    font_size=12.5)
add_para(tf2, "", font_size=4)
add_para(tf2,
    "4.  Consistent across both datasets (BGL & HDFS)  --  "
    "the finding is not dataset-specific.",
    font_size=12.5)

# ── Bottom-right: conclusion callout ─────────────────────────────
callout_bg = add_rounded_rect(8.5, 4.85, 4.4, 0.3, RGBColor(0x27, 0xAE, 0x60))
add_text_box(8.65, 4.85, 4.1, 0.3,
             "Conclusion: Train-only corpus is empirically validated. No leakage, no loss.",
             font_size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

# ── Save ──────────────────────────────────────────────────────────
out_path = "results/evidence_store_train_only_rationale.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
